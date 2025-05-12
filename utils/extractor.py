import fitz
from pptx import Presentation

def extract_text_from_pdf(path):
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc), len(doc)

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = ""
    count = 0
    for slide in prs.slides:
        count += 1
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text, count
