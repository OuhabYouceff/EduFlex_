import os
import json
import requests
from datetime import datetime
import pendulum
from flask import Flask, request, jsonify
from dotenv import load_dotenv
import fitz  # PyMuPDF for PDF processing
from pptx import Presentation  # python-pptx for PPTX processing

# Initialize Flask app
app = Flask(__name__)

# Load environment variables
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
GROQ_MODEL = os.getenv("GROQ_MODEL", "llama-3.2-8b-instruct")
GROQ_ENDPOINT = "https://api.groq.com/openai/v1/chat/completions"

def validate_time_slots(free_slots):
    """Validate the format and content of time slots."""
    try:
        if not isinstance(free_slots, list):
            raise ValueError("Free slots must be a list")
        for slot in free_slots:
            if not isinstance(slot, dict) or "start" not in slot or "end" not in slot:
                raise ValueError("Invalid time slot format")
            # Basic ISO 8601 validation
            start = pendulum.parse(slot["start"])
            end = pendulum.parse(slot["end"])
            if start >= end:
                raise ValueError("Start time must be before end time")
            # Ensure times are within allowed range (08:00 to 00:00)
            start_hour = start.hour + start.minute / 60
            end_hour = end.hour + end.minute / 60
            if start_hour < 8 or end_hour > 24:
                raise ValueError("Study sessions must be between 08:00 and 00:00")
        return True
    except Exception as e:
        print(f"Time slots validation error: {e}")
        return False

def normalize_time_slots(free_slots):
    """Normalize time slots to a consistent time zone (Europe/Paris, +02:00)."""
    try:
        normalized = []
        for slot in free_slots:
            start = pendulum.parse(slot["start"]).in_timezone("Europe/Paris")
            end = pendulum.parse(slot["end"]).in_timezone("Europe/Paris")
            normalized.append({
                "start": start.to_iso8601_string(),
                "end": end.to_iso8601_string()
            })
        return normalized
    except Exception as e:
        print(f"Time slots normalization error: {e}")
        return free_slots

def estimate_study_times_with_groq(filename, page_count):
    """Estimate study time based on file page count."""
    estimated_hours = max(1, round(page_count * 0.25))  # 4 pages ~ 1h
    return {
        "title": filename,
        "topics": [
            {
                "title": filename.replace(".pdf", "").replace(".pptx", "").replace("_", " "),
                "revision_time_hours": estimated_hours,
                "prerequisites": []
            }
        ]
    }

def generate_study_plan(curriculum, free_slots):
    """Generate a study plan based on curriculum and available time slots."""
    # Validate and normalize inputs
    if not validate_time_slots(free_slots):
        print("Invalid time slots provided")
        return None
    free_slots = normalize_time_slots(free_slots)

    # Debug input data
    print("📥 Curriculum:", json.dumps(curriculum, indent=2))
    print("📥 Free Slots:", json.dumps(free_slots, indent=2))

    prompt = f"""
You are a study planning assistant. Generate a detailed study plan based on the following curriculum and available time.

Constraints:
- Use ONLY the given available time slots.
- Never schedule study sessions before 08:00 or after 00:00.
- Max session: 2 hours.
- Add 15 minutes break between sessions.
- Between sessions put 10 minutes breks
- Start sessions at :00 or :30 only.
- Time format must be ISO 8601: "2025-04-06T10:00:00+02:00".

Curriculum:
{json.dumps(curriculum, indent=2)}

Available Time Slots:
{json.dumps(free_slots, indent=2)}

Output (JSON only):
[
  {{
    "course": "Course Title",
    "start": "2025-04-06T10:00:00+02:00",
    "end": "2025-04-06T12:00:00+02:00"
  }},
  {{
    "course": "Break",
    "start": "2025-04-06T10:00:00+02:00",
    "end": "2025-04-06T12:00:00+02:00"
  }},

]

Ensure the output is concise and contains only the JSON array, with no additional text.
"""

    headers = {
        "Authorization": f"Bearer {GROQ_API_KEY}",
        "Content-Type": "application/json"
    }

    payload = {
        "model": GROQ_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3,
        "max_tokens": 4096  # Increased to handle longer responses
    }

    try:
        response = requests.post(GROQ_ENDPOINT, headers=headers, json=payload)
        response.raise_for_status()
        content = response.json()["choices"][0]["message"]["content"]
        print("📤 LLM raw response:\n", content[:300])

        # Extract JSON array from response
        start = content.find("[")
        end = content.rfind("]") + 1
        if start == -1 or end == 0:
            print("❌ Invalid JSON: No array found in response")
            return None

        json_content = content[start:end]
        try:
            study_plan = json.loads(json_content)
            return study_plan
        except json.JSONDecodeError as e:
            print(f"❌ JSON Parse Error: {e}")
            print("📨 Full LLM response:", content)
            return None

    except requests.exceptions.RequestException as e:
        print(f"❌ API Request Error: {e}")
        return None
    except Exception as e:
        print(f"❌ LLM Study Plan Error: {e}")
        print("📨 Payload content (debug):", prompt[:1000])
        return None

def generate_quiz_from_file(file_path):
    """Generate a quiz from a PDF or PPTX file."""
    try:
        if file_path.endswith(".pdf"):
            doc = fitz.open(file_path)
            content = "\n".join(page.get_text() for page in doc)
        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            content = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        else:
            print(f"Unsupported file format: {file_path}")
            return []

        prompt = f"""
Generate a short quiz (max 5 questions) from the course below.
For each question, include a single correct answer.

Output format:
[
  {{
    "question": "What is...",
    "answer": "Correct answer"
  }},
  ...
]

COURSE:
{content[:3000]}
"""

        headers = {
            "Authorization": f"Bearer {GROQ_API_KEY}",
            "Content-Type": "application/json"
        }

        payload = {
            "model": GROQ_MODEL,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.4,
            "max_tokens": 1500
        }

        response = requests.post(GROQ_ENDPOINT, headers=headers, json=payload)
        response.raise_for_status()
        content = response.json()["choices"][0]["message"]["content"]
        start = content.find("[")
        end = content.rfind("]") + 1
        if start == -1 or end == 0:
            print("❌ Invalid JSON: No quiz array found")
            return []

        return json.loads(content[start:end])
    except Exception as e:
        print(f"❌ Quiz generation error: {e}")
        return []

def evaluate_quiz(user_answers):
    """Evaluate quiz answers (placeholder logic)."""
    correct = 0
    total = len(user_answers)
    for key, value in user_answers.items():
        if value.lower() == "a":  # Placeholder: assumes 'a' is correct
            correct += 1
    return correct / total if total > 0 else 0

@app.route("/planner", methods=["GET", "POST"])
def planner():
    """Handle requests to the /planner endpoint."""
    try:
        if request.method == "POST":
            data = request.get_json()
            if not data or "curriculum" not in data or "free_slots" not in data:
                return jsonify({"error": "Missing curriculum or free_slots"}), 400

            curriculum = data["curriculum"]
            free_slots = data["free_slots"]

            study_plan = generate_study_plan(curriculum, free_slots)
            if study_plan is None:
                return jsonify({"error": "Failed to generate study plan"}), 500

            return jsonify(study_plan), 200
        else:
            # For GET, return a sample response or form (adjust as needed)
            return jsonify({"message": "Send a POST request with curriculum and free_slots to generate a study plan"}), 200
    except Exception as e:
        print(f"❌ Planner endpoint error: {e}")
        return jsonify({"error": "Internal server error"}), 500

@app.route("/static/<path:path>")
def serve_static(path):
    """Serve static files (e.g., style.css)."""
    return app.send_static_file(path)

if __name__ == "__main__":
    app.run(host="127.0.0.1", port=5000, debug=True)